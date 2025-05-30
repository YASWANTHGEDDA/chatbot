# server/rag_service/document_parser.py

import os
import logging
from typing import List, Dict, Any, Optional
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import json
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - [%(name)s:%(lineno)d] - %(message)s')
logger = logging.getLogger(__name__)

class DocumentParser:
    """Parser for different document formats."""
    
    @staticmethod
    def parse_pdf(file_path: str) -> Dict[str, Any]:
        """Parse a PDF file and extract text and metadata."""
        try:
            logger.info(f"Parsing PDF file: {file_path}")
            
            # Read PDF
            reader = PdfReader(file_path)
            
            # Extract text from each page
            text_chunks = []
            for page in reader.pages:
                text = page.extract_text()
                if text.strip():
                    text_chunks.append(text)
            
            # Extract metadata
            metadata = reader.metadata
            if metadata:
                metadata = {
                    'title': metadata.get('/Title', ''),
                    'author': metadata.get('/Author', ''),
                    'subject': metadata.get('/Subject', ''),
                    'keywords': metadata.get('/Keywords', ''),
                    'creator': metadata.get('/Creator', ''),
                    'producer': metadata.get('/Producer', ''),
                    'creation_date': metadata.get('/CreationDate', ''),
                    'modification_date': metadata.get('/ModDate', '')
                }
            else:
                metadata = {}
            
            # Add file info to metadata
            metadata.update({
                'file_name': os.path.basename(file_path),
                'file_type': 'pdf',
                'page_count': len(reader.pages)
            })
            
            return {
                'text': '\n\n'.join(text_chunks),
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error parsing PDF file {file_path}: {e}", exc_info=True)
            raise
    
    @staticmethod
    def parse_docx(file_path: str) -> Dict[str, Any]:
        """Parse a DOCX file and extract text and metadata."""
        try:
            logger.info(f"Parsing DOCX file: {file_path}")
            
            # Read DOCX
            doc = DocxDocument(file_path)
            
            # Extract text from paragraphs
            text_chunks = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    text_chunks.append(text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            row_text.append(text)
                    if row_text:
                        text_chunks.append(' | '.join(row_text))
            
            # Extract metadata
            core_properties = doc.core_properties
            metadata = {
                'title': core_properties.title or '',
                'author': core_properties.author or '',
                'subject': core_properties.subject or '',
                'keywords': core_properties.keywords or '',
                'created': str(core_properties.created) if core_properties.created else '',
                'modified': str(core_properties.modified) if core_properties.modified else '',
                'file_name': os.path.basename(file_path),
                'file_type': 'docx'
            }
            
            return {
                'text': '\n\n'.join(text_chunks),
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error parsing DOCX file {file_path}: {e}", exc_info=True)
            raise
    
    @staticmethod
    def parse_pptx(file_path: str) -> Dict[str, Any]:
        """Parse a PPTX file and extract text and metadata."""
        try:
            logger.info(f"Parsing PPTX file: {file_path}")
            
            # Read PPTX
            prs = Presentation(file_path)
            
            # Extract text from slides
            text_chunks = []
            for slide in prs.slides:
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                
                if slide_text:
                    text_chunks.append(f"Slide {len(text_chunks) + 1}:\n" + '\n'.join(slide_text))
            
            # Extract metadata
            core_properties = prs.core_properties
            metadata = {
                'title': core_properties.title or '',
                'author': core_properties.author or '',
                'subject': core_properties.subject or '',
                'keywords': core_properties.keywords or '',
                'created': str(core_properties.created) if core_properties.created else '',
                'modified': str(core_properties.modified) if core_properties.modified else '',
                'file_name': os.path.basename(file_path),
                'file_type': 'pptx',
                'slide_count': len(prs.slides)
            }
            
            return {
                'text': '\n\n'.join(text_chunks),
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error parsing PPTX file {file_path}: {e}", exc_info=True)
            raise
    
    @staticmethod
    def parse_txt(file_path: str) -> Dict[str, Any]:
        """Parse a text file and extract content."""
        try:
            logger.info(f"Parsing text file: {file_path}")
            
            # Read text file
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            # Create metadata
            metadata = {
                'file_name': os.path.basename(file_path),
                'file_type': 'txt'
            }
            
            return {
                'text': text,
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error parsing text file {file_path}: {e}", exc_info=True)
            raise
    
    @staticmethod
    def parse_json(file_path: str) -> Dict[str, Any]:
        """Parse a JSON file and extract content."""
        try:
            logger.info(f"Parsing JSON file: {file_path}")
            
            # Read JSON file
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Convert JSON to formatted text
            text = json.dumps(data, indent=2)
            
            # Create metadata
            metadata = {
                'file_name': os.path.basename(file_path),
                'file_type': 'json'
            }
            
            return {
                'text': text,
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error parsing JSON file {file_path}: {e}", exc_info=True)
            raise
    
    @classmethod
    def parse_file(cls, file_path: str) -> Dict[str, Any]:
        """Parse a file based on its extension."""
        try:
            # Get file extension
            _, ext = os.path.splitext(file_path)
            ext = ext.lower()
            
            # Parse based on extension
            if ext == '.pdf':
                return cls.parse_pdf(file_path)
            elif ext == '.docx':
                return cls.parse_docx(file_path)
            elif ext == '.pptx':
                return cls.parse_pptx(file_path)
            elif ext == '.txt':
                return cls.parse_txt(file_path)
            elif ext == '.json':
                return cls.parse_json(file_path)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
            
        except Exception as e:
            logger.error(f"Error parsing file {file_path}: {e}", exc_info=True)
            raise 