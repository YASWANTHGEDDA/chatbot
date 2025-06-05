import logging
import os
from typing import List, Dict, Any, Optional, Tuple
from pypdf import PdfReader
from pptx import Presentation
from langchain_community.vectorstores import FAISS
from langchain_ollama import OllamaEmbeddings, ChatOllama
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from ..config.config import settings

logger = logging.getLogger(__name__)

class RAGService:
    def __init__(self):
        self.embeddings = None
        self.vector_store = None
        self.llm = None
        self.initialize_components()
        
    def initialize_components(self):
        """Initialize embeddings, vector store, and LLM."""
        try:
            # Initialize embeddings
            self.embeddings = OllamaEmbeddings(
                base_url=settings.OLLAMA_BASE_URL,
                model=settings.OLLAMA_EMBED_MODEL,
            )
            logger.info("Embeddings initialized successfully")
            
            # Initialize LLM for sub-query generation
            self.llm = ChatOllama(
                base_url=settings.OLLAMA_BASE_URL,
                model=settings.OLLAMA_MODEL,
                request_timeout=settings.OLLAMA_REQUEST_TIMEOUT
            )
            logger.info("LLM initialized successfully")
            
            # Load or initialize vector store
            self.load_vector_store()
            
        except Exception as e:
            logger.error(f"Error initializing RAG components: {e}", exc_info=True)
            raise
    
    def load_vector_store(self):
        """Load or initialize FAISS vector store."""
        try:
            if os.path.exists(settings.FAISS_FOLDER) and os.listdir(settings.FAISS_FOLDER):
                self.vector_store = FAISS.load_local(
                    settings.FAISS_FOLDER,
                    self.embeddings,
                    allow_dangerous_deserialization=True
                )
                logger.info(f"Loaded existing FAISS vector store from {settings.FAISS_FOLDER}")
            else:
                logger.info("No existing FAISS index found, will be created on first document addition")
                self.vector_store = None
                
        except Exception as e:
            logger.error(f"Error loading vector store: {e}", exc_info=True)
            raise
    
    def save_vector_store(self):
        """Save FAISS vector store."""
        try:
            if self.vector_store:
                self.vector_store.save_local(settings.FAISS_FOLDER)
                logger.info(f"Saved FAISS vector store to {settings.FAISS_FOLDER}")
        except Exception as e:
            logger.error(f"Error saving vector store: {e}", exc_info=True)
            raise
    
    def extract_text_from_file(self, filepath: str) -> str:
        """Extract text from PDF or PPT file."""
        try:
            ext = os.path.splitext(filepath)[1].lower()
            if ext == '.pdf':
                reader = PdfReader(filepath)
                if not reader.pages:
                    logger.error(f"PDF file {filepath} has no pages or could not be read")
                    return ""
                text = ""
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                    else:
                        logger.warning(f"No text extracted from page {i+1} of {filepath}")
                return text.strip() or ""
                
            elif ext in ('.pptx', '.ppt'):
                prs = Presentation(filepath)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text.strip() or ""
                
            else:
                logger.warning(f"Unsupported file extension for text extraction: {ext} in {filepath}")
                return ""
                
        except Exception as e:
            logger.error(f"Error extracting text from '{filepath}': {e}", exc_info=True)
            return ""
    
    def create_chunks_from_text(self, text: str, filename: str) -> List[Document]:
        """Create document chunks from text."""
        try:
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=settings.RAG_CHUNK_SIZE,
                chunk_overlap=settings.RAG_CHUNK_OVERLAP
            )
            chunks = splitter.split_text(text)
            return [Document(page_content=chunk, metadata={"source": filename}) for chunk in chunks]
        except Exception as e:
            logger.error(f"Error creating chunks for '{filename}': {e}", exc_info=True)
            return []
    
    def add_documents_to_vector_store(self, documents: List[Document]) -> bool:
        """Add documents to FAISS vector store."""
        if not self.embeddings:
            logger.error("Cannot add documents to vector store: Embeddings not initialized")
            return False
            
        try:
            if not self.vector_store:
                logger.info("Creating new FAISS vector store")
                self.vector_store = FAISS.from_documents(documents, self.embeddings)
            else:
                self.vector_store.add_documents(documents)
            self.save_vector_store()
            return True
        except Exception as e:
            logger.error(f"Error adding documents to vector store: {e}", exc_info=True)
            return False
    
    async def generate_sub_queries(self, query: str) -> List[str]:
        """Generate sub-queries for RAG."""
        if not self.llm:
            logger.error("LLM not initialized, cannot generate sub-queries")
            return []
            
        try:
            chain = LLMChain(
                llm=self.llm,
                prompt=PromptTemplate(
                    template=settings.SUB_QUERY_PROMPT_TEMPLATE,
                    input_variables=["query", "num_queries"]
                )
            )
            response = await chain.arun(
                query=query,
                num_queries=settings.MULTI_QUERY_COUNT
            )
            return [q.strip() for q in response.split('\n') if q.strip()]
        except Exception as e:
            logger.error(f"Error generating sub-queries: {e}", exc_info=True)
            return []
    
    async def perform_rag_search(
        self,
        query: str,
        include_sub_queries: bool = True
    ) -> Tuple[List[Document], str, Dict[int, Dict[str, str]]]:
        """
        Perform RAG search with optional sub-query generation.
        
        Returns:
            Tuple of (retrieved_documents, context_text, docs_map)
        """
        if not self.vector_store:
            logger.warning("Vector store not loaded/initialized. RAG search cannot be performed")
            return [], "Knowledge base is currently unavailable.", {}
            
        try:
            # Generate sub-queries if enabled
            sub_queries = await self.generate_sub_queries(query) if include_sub_queries else []
            all_queries = [query] + sub_queries
            
            # Perform search for each query
            all_retrieved_docs_with_scores = []
            for q_text in all_queries:
                results_with_scores = self.vector_store.similarity_search_with_score(
                    q_text,
                    k=settings.RAG_SEARCH_K_PER_QUERY
                )
                all_retrieved_docs_with_scores.extend(results_with_scores)
            
            # De-duplicate documents based on content and source
            unique_docs_dict = {}
            for doc, score in all_retrieved_docs_with_scores:
                doc_key = (doc.metadata['source'], doc.page_content)
                if doc_key not in unique_docs_dict or score < unique_docs_dict[doc_key][1]:
                    unique_docs_dict[doc_key] = (doc, score)
            
            # Sort unique documents by score and take top RAG_CHUNK_K
            sorted_unique_docs = sorted(list(unique_docs_dict.values()), key=lambda item: item[1])
            final_docs = [item[0] for item in sorted_unique_docs[:settings.RAG_CHUNK_K]]
            
            if not final_docs:
                logger.info(f"No relevant documents found for query: {query}")
                return [], "No specific documents found for your query.", {}
            
            # Format context and create docs map
            context_text = "\n\n".join(
                f"[Chunk {i+1} from: {doc.metadata['source']}]\n{doc.page_content}"
                for i, doc in enumerate(final_docs)
            )
            docs_map = {
                i+1: {
                    "source": doc.metadata['source'],
                    "content": doc.page_content
                }
                for i, doc in enumerate(final_docs)
            }
            
            return final_docs, context_text, docs_map
            
        except Exception as e:
            logger.error(f"Error in RAG search: {e}", exc_info=True)
            return [], "Error retrieving documents.", {}
    
    async def process_document(self, filepath: str) -> bool:
        """Process a document and add it to the vector store."""
        try:
            # Extract text
            text = self.extract_text_from_file(filepath)
            if not text:
                return False
            
            # Create chunks
            filename = os.path.basename(filepath)
            chunks = self.create_chunks_from_text(text, filename)
            if not chunks:
                return False
            
            # Add to vector store
            return self.add_documents_to_vector_store(chunks)
            
        except Exception as e:
            logger.error(f"Error processing document '{filepath}': {e}", exc_info=True)
            return False 