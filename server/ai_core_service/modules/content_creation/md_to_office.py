# modules/content_creation/md_to_office.py
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DocxInches # Use a different alias for docx Inches

def parse_slide_chunk(chunk):
    """Parses a single slide's text block to extract its components."""
    slide_title = "Untitled Slide"
    
    # Match "### Slide X: Title" or "### Title"
    title_match = re.match(r"### (?:Slide \d+:\s*)?(.*)", chunk)
    if title_match:
        slide_title = title_match.group(1).strip()

    text_content_match = re.search(
        r"\*\*Slide Text Content:\*\*(.*?)(?=\*\*Image Prompt:\*\*|\*\*Author Notes for Slide|\Z)", 
        chunk, 
        re.DOTALL
    )
    raw_text_content = text_content_match.group(1).strip() if text_content_match else ""

    image_prompt_match = re.search(
        r"\*\*Image Prompt:\*\*(.*?)(?=\n---\n\*\*Author Notes for Slide|\*\*Author Notes for Slide|\Z)",
        chunk,
        re.DOTALL
    )
    raw_image_prompt = image_prompt_match.group(1).strip() if image_prompt_match else ""
    
    author_notes_match = re.search(
        r"\*\*Author Notes for Slide \d+:\*\*(.*)", 
        chunk, 
        re.DOTALL
    )
    raw_author_notes = author_notes_match.group(1).strip() if author_notes_match else ""
    
    return {
        "title": slide_title,
        "text_content": raw_text_content,
        "image_prompt": raw_image_prompt,
        "author_notes": raw_author_notes
    }

def refined_parse_markdown(markdown_content):
    """Splits the markdown content into slide chunks and parses each chunk."""
    slides_data = []
    # Find the start of the first slide definition "### Slide X:" or "### Some Title"
    # This handles cases where there might be introductory text before the first slide.
    first_slide_match = re.search(r"### (?:Slide \d+:|[^-\n][^\n]*)", markdown_content)
    
    processed_slides_content = markdown_content
    if first_slide_match:
        processed_slides_content = markdown_content[first_slide_match.start():]
    else: # No identifiable slides
        return [] # Return empty list if no slides are found
    
    # Split by "---" separator, but only if it's followed by "### Slide X:" or end of string
    # This regex uses a positive lookahead
    slide_chunks = re.split(r'\n\s*---\s*\n(?=### (?:Slide \d+:|[^-\n][^\n]*))', processed_slides_content)

    for chunk_idx, chunk in enumerate(slide_chunks):
        stripped_chunk = chunk.strip()
        # A chunk is valid if it starts with "### " (indicating a title)
        if not stripped_chunk or not stripped_chunk.startswith("### "):
            continue
        
        data = parse_slide_chunk(stripped_chunk)
        # Only add if there's some content
        if data["title"] != "Untitled Slide" or data["text_content"] or data["image_prompt"] or data["author_notes"]:
            slides_data.append(data)
            
    return slides_data

def add_text_to_shape_with_markdown(text_frame, markdown_text, is_title=False):
    """Adds markdown-formatted text to a PowerPoint text frame."""
    text_frame.clear() 
    text_frame.word_wrap = True

    title_font_size = Pt(30)
    content_font_size = Pt(16) # Default content font size

    lines = markdown_text.split('\n')
    if not lines and not markdown_text.strip(): # Handle completely empty text_content
        p = text_frame.add_paragraph()
        run = p.add_run(); run.text = " " # Add a space to make the textbox not collapse if it's the only one
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.size = content_font_size if not is_title else title_font_size
        return

    for line_num, line in enumerate(lines):
        p = text_frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        
        bullet_match = re.match(r'^(\s*)([\*\-])\s*(.*)', line) 
        
        is_bullet = False
        content_line = line 
        indent_level = 0

        if bullet_match:
            leading_spaces = bullet_match.group(1)
            indent_level = min(len(leading_spaces) // 2, 8) # Assuming 2 spaces per indent level
            content_line = bullet_match.group(3) 
            is_bullet = True
            p.level = indent_level 
        else:
            content_line = line.lstrip() 

        segments = re.split(r'(\*\*.*?\*\*|__.*?__)', content_line)
        
        has_actual_text_in_segments = any(s.strip() for s in segments)

        if not has_actual_text_in_segments and line.strip() == "": # Preserve blank lines
             # An empty paragraph 'p' is already added, which creates a blank line.
            if is_bullet: # if it was like "*  " (bullet with only spaces after)
                run = p.add_run(); run.text = " "
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.size = content_font_size

        elif not has_actual_text_in_segments and is_bullet: # e.g. "* "
            run = p.add_run(); run.text = " "
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.size = content_font_size
        
        else: # Has actual text content
            for segment_idx, segment in enumerate(segments):
                if not segment and segment_idx == 0 and not content_line.startswith(("*","__")):
                    # Handles cases where a line starts with non-bold/italic text
                    # and re.split might produce an initial empty segment.
                    # This part is tricky; simpler to ensure segments are not empty.
                    pass

                if not segment: continue 
                
                run = p.add_run()
                if (segment.startswith("**") and segment.endswith("**")) or \
                   (segment.startswith("__") and segment.endswith("__")):
                    run.text = segment[2:-2]
                    run.font.bold = True
                else:
                    run.text = segment
                
                run.font.color.rgb = RGBColor(255, 255, 255) # White text
                if is_title:
                    run.font.size = title_font_size
                else:
                    run.font.size = content_font_size

def create_ppt(slides_data, output_dir, filename="Presentation.pptx"):
    """
    Creates a PowerPoint presentation from parsed slide data.

    Args:
        slides_data (list): List of slide data dictionaries.
        output_dir (str): Directory to save the PPTX file.
        filename (str): Name of the PPTX file.

    Returns:
        str: Full path to the created PPTX file, or None if failed.
    """
    os.makedirs(output_dir, exist_ok=True)
    output_filepath = os.path.join(output_dir, filename)

    prs = Presentation()
    prs.slide_width = Inches(16) 
    prs.slide_height = Inches(9)

    if not slides_data:
        slide_layout = prs.slide_layouts[5] # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background; fill = background.fill
        fill.solid(); fill.fore_color.rgb = RGBColor(0, 0, 0) # Black background
        
        # Add a title shape for the "no content" message
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
        tf = title_shape.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "No slide content found in input."
        run.font.color.rgb = RGBColor(255, 255, 255); run.font.size = Pt(24)
        
        prs.save(output_filepath)
        print(f"PPTX file '{output_filepath}' created (empty content due to no slide data).")
        return output_filepath

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[5] # Use a blank slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Set background color (black)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0) # Black

        # Add title textbox
        # Adjust position and size as needed (Left, Top, Width, Height)
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), prs.slide_width - Inches(1.0), Inches(1.0))
        add_text_to_shape_with_markdown(title_shape.text_frame, slide_data["title"], is_title=True)
        
        # Add content textbox if text_content is present
        if slide_data["text_content"]:
            # Position content to the left, leaving space for a potential image placeholder on the right
            content_shape = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), 
                (prs.slide_width * 0.65) - Inches(0.5), # ~65% width for text
                prs.slide_height - Inches(2.0) # Adjust height
            )
            add_text_to_shape_with_markdown(content_shape.text_frame, slide_data["text_content"])
        
        # Note: Image generation/insertion from "image_prompt" is not handled here.
        # That would require an image generation API and `slide.shapes.add_picture()`.

    prs.save(output_filepath)
    print(f"PPTX file '{output_filepath}' created successfully with {len(slides_data)} slides.")
    return output_filepath

def add_markdown_line_to_docx(doc_parent, markdown_line):
    """Adds a single line of markdown-processed text to a docx document/parent (doc or table cell)."""
    bullet_match = re.match(r'^(\s*)([\*\-])\s*(.*)', markdown_line)

    if bullet_match:
        leading_spaces = bullet_match.group(1)
        content_line = bullet_match.group(3)
        indent_level = min(len(leading_spaces) // 2, 8) 

        p_style = 'ListBullet'
        # python-docx doesn't directly support multiple list levels with one style easily.
        # For deeper indents, you might need to define ListBullet2, ListBullet3, etc., in your template
        # or manually adjust paragraph_format.left_indent.
        if indent_level > 0:
             p_style = 'ListBullet' # Still use ListBullet, but adjust indent

        p = doc_parent.add_paragraph(style=p_style)
        if indent_level > 0:
            p.paragraph_format.left_indent = DocxInches(0.25 * indent_level) # 0.25 inch per level

        segments = re.split(r'(\*\*.*?\*\*|__.*?__)', content_line)
        added_text_run = False
        for segment in segments:
            if not segment: continue
            run = p.add_run()
            if (segment.startswith("**") and segment.endswith("**")) or \
               (segment.startswith("__") and segment.endswith("__")):
                run.text = segment[2:-2]
                run.font.bold = True
            else:
                run.text = segment
            if segment.strip(): # check if the segment adds actual text
                added_text_run = True
        
        if not added_text_run and not content_line.strip(): # If bullet item is empty (e.g. "* ")
            run = p.add_run()
            run.text = " " # Add a space to make the bullet visible
    
    else: # Not a bullet point
        stripped_line = markdown_line.lstrip() # Remove leading spaces from non-bullet lines
        
        if not stripped_line and markdown_line == "": # An actual empty line from markdown
            doc_parent.add_paragraph("") # Preserve blank line
            return
        elif not stripped_line: # Line was only whitespace, treat as blank
            doc_parent.add_paragraph("")
            return

        p = doc_parent.add_paragraph()
        segments = re.split(r'(\*\*.*?\*\*|__.*?__)', stripped_line)
        for segment in segments:
            if not segment: continue
            run = p.add_run()
            if (segment.startswith("**") and segment.endswith("**")) or \
               (segment.startswith("__") and segment.endswith("__")):
                run.text = segment[2:-2]
                run.font.bold = True
            else:
                run.text = segment

def create_doc(slides_data, output_dir, doc_filename, content_key):
    """
    Creates a Word document from parsed slide data for a specific content key.

    Args:
        slides_data (list): List of slide data dictionaries.
        output_dir (str): Directory to save the DOCX file.
        doc_filename (str): Name of the DOCX file.
        content_key (str): The key in slide_data dict to extract content from (e.g., "image_prompt", "author_notes").

    Returns:
        str: Full path to the created DOCX file, or None if failed.
    """
    os.makedirs(output_dir, exist_ok=True)
    output_filepath = os.path.join(output_dir, doc_filename)
    
    doc = Document()
    # Use the filename (without extension) as the main document title
    doc_title = os.path.splitext(doc_filename)[0]
    doc.add_heading(doc_title, level=0) # Level 0 for main document title
    
    if not slides_data:
        doc.add_paragraph(f"[No slide data provided to extract for '{content_key}']")
        doc.save(output_filepath)
        print(f"DOCX file '{output_filepath}' created (empty content due to no slide data).")
        return output_filepath

    for i, slide_data in enumerate(slides_data):
        doc.add_heading(f"Slide: {slide_data.get('title', 'Untitled Slide')}", level=1) # Level 1 for slide titles
        
        content_to_add = slide_data.get(content_key, "")

        if not content_to_add or not content_to_add.strip():
            doc.add_paragraph(f"[No content for '{content_key}' in this slide]")
        else:
            lines = content_to_add.split('\n')
            for line_text in lines:
                add_markdown_line_to_docx(doc, line_text) # Pass doc object as parent

        if i < len(slides_data) - 1:
            # Add a more distinct visual separator if desired, e.g., a horizontal rule
            # doc.add_page_break() # Or a page break
            doc.add_paragraph("\n" + "_" * 30 + "\n") # Simple text separator

    doc.save(output_filepath)
    print(f"DOCX file '{output_filepath}' created successfully with content from '{content_key}'.")
    return output_filepath