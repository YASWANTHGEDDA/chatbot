# --- START OF FILE ai_core.py ---

import os
import logging
import json
import ollama
import concurrent.futures
from tqdm import tqdm
from pypdf import PdfReader
from pptx import Presentation
from langchain_community.vectorstores import FAISS
from langchain_ollama import OllamaEmbeddings, ChatOllama
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from langchain.chains import LLMChain
# PromptTemplate is already imported in config, but good practice if used directly here too
# from langchain.prompts import PromptTemplate

from config import (
    OLLAMA_BASE_URL, OLLAMA_MODEL, OLLAMA_EMBED_MODEL, FAISS_FOLDER,
    DEFAULT_PDFS_FOLDER, UPLOAD_FOLDER, RAG_CHUNK_K, MULTI_QUERY_COUNT,
    ANALYSIS_MAX_CONTEXT_LENGTH, OLLAMA_REQUEST_TIMEOUT, RAG_SEARCH_K_PER_QUERY,
    SUB_QUERY_PROMPT_TEMPLATE, SYNTHESIS_PROMPT_TEMPLATE, ANALYSIS_PROMPTS,
    KG_CHUNK_SIZE, KG_CHUNK_OVERLAP, KG_MAX_WORKERS, KG_MODEL,
    KG_OUTPUT_FOLDER, KG_PROMPT_TEMPLATE, KG_FILENAME_SUFFIX # Added KG_FILENAME_SUFFIX
)
from utils import parse_llm_response, escape_html # Assuming you have these

logger = logging.getLogger(__name__)

# Global State
document_texts_cache = {}
vector_store = None
embeddings = None
llm = None
_kg_ollama_client = None

# Initialization Functions
def initialize_ai_components():
    """Initializes embeddings and LLM."""
    global embeddings, llm
    logger.info("Initializing AI components...")
    try:
        embeddings = OllamaEmbeddings(
            base_url=OLLAMA_BASE_URL,
            model=OLLAMA_EMBED_MODEL,
        )
        llm = ChatOllama(
            base_url=OLLAMA_BASE_URL,
            model=OLLAMA_MODEL,
            request_timeout=OLLAMA_REQUEST_TIMEOUT
        )
        logger.info("AI components initialized.")
        return embeddings, llm
    except Exception as e:
        logger.error(f"Failed to initialize AI components: {e}", exc_info=True)
        return None, None

def load_vector_store():
    """Loads or initializes FAISS vector store."""
    global vector_store
    try:
        if os.path.exists(FAISS_FOLDER) and os.listdir(FAISS_FOLDER): # Check if directory is not empty
            vector_store = FAISS.load_local(FAISS_FOLDER, embeddings, allow_dangerous_deserialization=True)
            logger.info(f"Loaded existing FAISS vector store from {FAISS_FOLDER}.")
        else:
            logger.info(f"No FAISS index found at {FAISS_FOLDER} or it's empty. Starting with empty index.")
            # Optionally initialize an empty store if needed immediately,
            # or let it be created on first add_documents_to_vector_store
        return True
    except Exception as e:
        logger.error(f"Failed to load FAISS vector store: {e}", exc_info=True)
        return False

def save_vector_store():
    """Saves FAISS vector store."""
    try:
        if vector_store:
            vector_store.save_local(FAISS_FOLDER)
            logger.info(f"Saved FAISS vector store to {FAISS_FOLDER}.")
    except Exception as e:
        logger.error(f"Failed to save FAISS vector store: {e}", exc_info=True)

def load_all_document_texts():
    """Loads text from all documents into cache."""
    import config # Explicit import
    logger.info("Loading all document texts into cache...")
    loaded_count = 0
    for folder_path in [config.DEFAULT_PDFS_FOLDER, config.UPLOAD_FOLDER]:
        if not os.path.exists(folder_path):
            logger.warning(f"Document folder not found: {folder_path}")
            continue
        for filename in os.listdir(folder_path):
            if filename.lower().endswith(tuple(config.ALLOWED_EXTENSIONS)) and not filename.startswith('~'):
                filepath = os.path.join(folder_path, filename)
                if filename not in document_texts_cache: # Avoid reloading if already cached
                    text = extract_text_from_file(filepath) # Renamed for clarity
                    if text:
                        document_texts_cache[filename] = text
                        loaded_count +=1
    logger.info(f"Loaded texts for {loaded_count} new documents into cache. Total cached: {len(document_texts_cache)}")


# Text Extraction and Chunking
def extract_text_from_file(filepath: str) -> str: # Renamed from extract_text_from_pdf
    """Extracts text from PDF or PPT file."""
    try:
        ext = os.path.splitext(filepath)[1].lower()
        if ext == '.pdf':
            reader = PdfReader(filepath)
            if not reader.pages:
                logger.error(f"PDF file {filepath} has no pages or could not be read.")
                return ""
            text = ""
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
                else:
                    logger.warning(f"No text extracted from page {i+1} of {filepath}.")
            return text.strip() or ""
        elif ext in ('.pptx', '.ppt'):
            prs = Presentation(filepath)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip() or ""
        else:
            logger.warning(f"Unsupported file extension for text extraction: {ext} in {filepath}")
            return ""
    except FileNotFoundError:
        logger.error(f"File not found for text extraction: {filepath}")
        return ""
    except Exception as e:
        logger.error(f"Error extracting text from '{filepath}': {e}", exc_info=True)
        return ""

def create_chunks_from_text(text: str, filename: str) -> list[Document]:
    """Creates document chunks from text."""
    try:
        # TODO: Consider making chunk_size and chunk_overlap configurable from config.py for RAG
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        chunks = splitter.split_text(text)
        return [Document(page_content=chunk, metadata={"source": filename}) for chunk in chunks]
    except Exception as e:
        logger.error(f"Error creating chunks for '{filename}': {e}", exc_info=True)
        return []

def add_documents_to_vector_store(documents: list[Document]) -> bool:
    """Adds documents to FAISS vector store."""
    global vector_store
    if not embeddings:
        logger.error("Cannot add documents to vector store: Embeddings not initialized.")
        return False
    try:
        if not vector_store:
            logger.info("Creating new FAISS vector store.")
            vector_store = FAISS.from_documents(documents, embeddings)
        else:
            vector_store.add_documents(documents)
        save_vector_store()
        return True
    except Exception as e:
        logger.error(f"Error adding documents to vector store: {e}", exc_info=True)
        return False

# RAG and Analysis
def generate_sub_queries(query: str) -> list[str]:
    """Generates sub-queries for RAG."""
    if not llm:
        logger.error("LLM not initialized, cannot generate sub-queries.")
        return []
    try:
        chain = LLMChain(llm=llm, prompt=SUB_QUERY_PROMPT_TEMPLATE)
        response = chain.run(query=query, num_queries=MULTI_QUERY_COUNT)
        return [q.strip() for q in response.split('\n') if q.strip()]
    except Exception as e:
        logger.error(f"Error generating sub-queries: {e}", exc_info=True)
        return []

def perform_rag_search(query: str) -> tuple[list[Document], str, dict]:
    """Performs RAG search."""
    if not vector_store:
        logger.warning("Vector store not loaded/initialized. RAG search cannot be performed.")
        return [], "Knowledge base is currently unavailable.", {}
    try:
        sub_queries = generate_sub_queries(query) if MULTI_QUERY_COUNT > 0 else []
        all_queries = [query] + sub_queries
        
        all_retrieved_docs_with_scores = []
        for q_idx, q_text in enumerate(all_queries):
            # Using similarity_search_with_score to potentially de-duplicate or rank later
            results_with_scores = vector_store.similarity_search_with_score(q_text, k=RAG_SEARCH_K_PER_QUERY)
            all_retrieved_docs_with_scores.extend(results_with_scores)

        # De-duplicate documents based on content and source, keeping the one with the best score (lower is better for L2/cosine)
        unique_docs_dict = {}
        for doc, score in all_retrieved_docs_with_scores:
            doc_key = (doc.metadata['source'], doc.page_content)
            if doc_key not in unique_docs_dict or score < unique_docs_dict[doc_key][1]:
                unique_docs_dict[doc_key] = (doc, score)
        
        # Sort unique documents by score and take top RAG_CHUNK_K
        sorted_unique_docs = sorted(list(unique_docs_dict.values()), key=lambda item: item[1])
        final_docs = [item[0] for item in sorted_unique_docs[:RAG_CHUNK_K]]

        if not final_docs:
            logger.info(f"No relevant documents found for query: {query}")
            return [], "No specific documents found for your query.", {}

        context_text = "\n\n".join(f"[Chunk {i+1} from: {doc.metadata['source']}]\n{doc.page_content}" for i, doc in enumerate(final_docs))
        # Map for citation: Key is citation number (1-based), value is details
        docs_map = {i+1: {"source": doc.metadata['source'], "content": doc.page_content} for i, doc in enumerate(final_docs)}
        
        return final_docs, context_text, docs_map
    except Exception as e:
        logger.error(f"Error in RAG search: {e}", exc_info=True)
        return [], "Error retrieving documents.", {}

# --- KG Related Functions ---
def get_kg_filepath(doc_filename: str) -> str:
    """Gets the expected filepath for a document's KG."""
    base_name = os.path.splitext(doc_filename)[0]
    kg_file = f"{base_name}{KG_FILENAME_SUFFIX}"
    return os.path.join(KG_OUTPUT_FOLDER, kg_file)

def load_knowledge_graph(doc_filename: str) -> dict | None:
    """Loads the knowledge graph for a given document filename."""
    kg_filepath = get_kg_filepath(doc_filename)
    if os.path.exists(kg_filepath):
        try:
            with open(kg_filepath, 'r', encoding='utf-8') as f:
                kg_data = json.load(f)
            logger.info(f"Loaded KG for '{doc_filename}' from {kg_filepath}")
            return kg_data
        except Exception as e:
            logger.error(f"Error loading KG JSON from {kg_filepath}: {e}", exc_info=True)
    else:
        logger.debug(f"No KG found for '{doc_filename}' at {kg_filepath}")
    return None

def format_kg_info_for_llm(kg_data: dict, query: str, max_insights: int = 3, max_details_per_node: int = 2) -> str:
    """
    Extracts and formats relevant KG information based on the query for LLM context.
    """
    if not kg_data or not isinstance(kg_data, dict) or 'nodes' not in kg_data:
        return ""

    query_keywords = set(q.lower() for q in query.split() if len(q) > 3) # Keywords from query
    
    relevant_node_infos = []
    nodes_by_id = {node.get('id'): node for node in kg_data.get('nodes', []) if isinstance(node, dict) and node.get('id')}

    # Find nodes relevant to the query
    for node_id, node_data in nodes_by_id.items():
        description = node_data.get('description', '').lower()
        node_id_lower = str(node_id).lower()
        if any(keyword in description for keyword in query_keywords) or \
           any(keyword in node_id_lower for keyword in query_keywords):
            relevant_node_infos.append(node_data)
    
    if not relevant_node_infos:
        return ""

    output_str = "Key concepts and relationships from Knowledge Graph:\n"
    insights_count = 0

    for node in relevant_node_infos:
        if insights_count >= max_insights:
            break
        
        node_id = node.get('id')
        desc = node.get('description', 'N/A')
        node_type = node.get('type', 'Concept')
        output_str += f"- Node: '{node_id}' (Type: {node_type})\n  Description: {desc}\n"
        
        # Find connected edges
        related_count = 0
        for edge in kg_data.get('edges', []):
            if related_count >= max_details_per_node: break
            if not isinstance(edge, dict): continue
            
            rel_from = edge.get('from')
            rel_to = edge.get('to')
            relationship = edge.get('relationship')

            if rel_from == node_id and rel_to in nodes_by_id:
                output_str += f"  - Relates to '{rel_to}' (via: {relationship})\n"
                related_count +=1
            elif rel_to == node_id and rel_from in nodes_by_id:
                output_str += f"  - Is related from '{rel_from}' (via: {relationship})\n"
                related_count +=1
        insights_count += 1
    
    return output_str.strip() if insights_count > 0 else ""

# --- End KG Related Functions ---


def synthesize_chat_response(query: str, context_text: str, context_docs_map: dict) -> tuple[str, str]:
    """Synthesizes chat response using RAG context and KG insights."""
    if not llm:
        logger.error("LLM not initialized, cannot synthesize response.")
        return "Error: AI model is not available.", ""
    try:
        # --- KG Integration ---
        kg_derived_insights = ""
        unique_source_files = set()
        if context_docs_map: # Ensure there's RAG context to derive KG from
            for doc_detail in context_docs_map.values():
                unique_source_files.add(doc_detail['source'])
        
        if unique_source_files:
            all_kg_texts = []
            for source_file in unique_source_files:
                kg_data = load_knowledge_graph(source_file)
                if kg_data:
                    formatted_kg_text = format_kg_info_for_llm(kg_data, query)
                    if formatted_kg_text:
                         all_kg_texts.append(f"For document '{source_file}':\n{formatted_kg_text}")
            if all_kg_texts:
                kg_derived_insights = "\n\n--- Knowledge Graph Insights ---\n" + "\n\n".join(all_kg_texts)
        # --- End KG Integration ---

        # Combine RAG context with KG insights
        enriched_context = context_text + kg_derived_insights
        logger.debug(f"Enriched context for LLM (len {len(enriched_context)}):\n{enriched_context[:1000]}...")


        chain = LLMChain(llm=llm, prompt=SYNTHESIS_PROMPT_TEMPLATE)
        # The prompt expects 'context' as the variable name for all contextual info
        response_text = chain.run(query=query, context=enriched_context) 
        
        thinking = response_text.split('</thinking>')[0].replace('<thinking>', '').strip() if '</thinking>' in response_text else ""
        answer = response_text.split('</thinking>')[-1].strip()
        return answer, thinking
    except Exception as e:
        logger.error(f"Error synthesizing chat response: {e}", exc_info=True)
        return f"Error: Could not generate a response. {str(e)}", ""


def generate_document_analysis(filename: str, analysis_type: str) -> tuple[str, str]:
    """Generates document analysis."""
    import config # Explicit import
    if not llm:
        logger.error("LLM not initialized, cannot generate document analysis.")
        return "Error: AI model is not available.", ""
    try:
        text = document_texts_cache.get(filename)
        if not text:
            # Attempt to load from disk if not in cache
            file_path_to_try = os.path.join(UPLOAD_FOLDER, filename)
            if not os.path.exists(file_path_to_try):
                file_path_to_try = os.path.join(DEFAULT_PDFS_FOLDER, filename)
            
            if os.path.exists(file_path_to_try):
                logger.info(f"Text for '{filename}' not in cache, loading from {file_path_to_try} for analysis.")
                text = extract_text_from_file(file_path_to_try)
                if text:
                    document_texts_cache[filename] = text # Cache it now
            else:
                logger.error(f"File '{filename}' not found for analysis.")
                return None, "File not found."

        if not text:
            return None, "No text available for analysis."

        chain = LLMChain(llm=llm, prompt=config.ANALYSIS_PROMPTS[analysis_type])
        # Ensure text is not excessively long for this type of analysis
        response_text = chain.run(doc_text_for_llm=text[:ANALYSIS_MAX_CONTEXT_LENGTH]) 
        thinking = response_text.split('</thinking>')[0].replace('<thinking>', '').strip() if '</thinking>' in response_text else ""
        content = response_text.split('</thinking>')[-1].strip()
        return content, thinking
    except Exception as e:
        logger.error(f"Error in document analysis for '{filename}': {e}", exc_info=True)
        return f"Error: {str(e)}", ""

# Knowledge Graph Generation (largely same, but saving part needs update)
def _initialize_kg_ollama_client() -> ollama.Client | None:
    global _kg_ollama_client
    if _kg_ollama_client:
        return _kg_ollama_client
    try:
        logger.info(f"Initializing KG Ollama client: {OLLAMA_BASE_URL}, model={KG_MODEL}")
        _kg_ollama_client = ollama.Client(host=OLLAMA_BASE_URL)
        _kg_ollama_client.list() # Test connection
        logger.info("KG Ollama client initialized.")
        return _kg_ollama_client
    except Exception as e:
        logger.error(f"Failed to initialize KG Ollama client: {e}", exc_info=True)
        _kg_ollama_client = None
        return None

def _kg_split_into_chunks(text: str, chunk_size: int, overlap: int) -> list[str]:
    logger.info(f"KG: Splitting text into chunks (size={chunk_size}, overlap={overlap})...")
    chunks = []
    start = 0
    text_len = len(text)
    while start < text_len:
        end = min(start + chunk_size, text_len)
        chunks.append(text[start:end])
        next_start = end - overlap
        # Ensure progress, especially with small texts or large overlaps
        if next_start <= start and end < text_len :
            start = start + 1 
        else:
            start = next_start if end < text_len else end
    logger.info(f"KG: Split into {len(chunks)} chunks.")
    return chunks


def _kg_process_single_chunk(chunk_data: tuple[int, str], ollama_client_instance: ollama.Client, model_name: str, prompt_template_str: str) -> dict | None:
    index, chunk_text = chunk_data
    chunk_num = index + 1
    
    # The KG_PROMPT_TEMPLATE from config.py is already escaped correctly
    full_prompt = prompt_template_str.format(chunk_text=chunk_text) 
    
    if not ollama_client_instance:
        logger.error(f"KG: Ollama client not available for chunk {chunk_num}.")
        return None
    try:
        response = ollama_client_instance.chat(
            model=model_name,
            messages=[{"role": "user", "content": full_prompt}],
            format="json", # Ollama should handle this based on model capabilities
            options={"num_ctx": 4096, "temperature": 0.3} # Example options
        )
        content = response.get('message', {}).get('content', '')
        if not content:
            logger.warning(f"KG: Empty response for chunk {chunk_num}")
            return {}
        
        # Basic cleanup for ```json ... ``` markdown blocks if model adds them
        if content.strip().startswith('```json'):
            content = content.strip()[7:-3].strip()
        elif content.strip().startswith('```'): # More generic ``` removal
            content = content.strip()[3:-3].strip()
            
        graph_data = json.loads(content) # Expecting JSON directly now
        
        # Validate basic structure
        if isinstance(graph_data, dict) and \
           'nodes' in graph_data and isinstance(graph_data['nodes'], list) and \
           'edges' in graph_data and isinstance(graph_data['edges'], list):
            return graph_data
        
        logger.warning(f"KG: Invalid graph structure for chunk {chunk_num}. Content: {content[:200]}...")
        return {}
    except json.JSONDecodeError as je:
        logger.error(f"KG: JSON Decode Error processing chunk {chunk_num}: {je}. Content: {content[:500]}...", exc_info=False) # Don't need full exc_info for common JSON error
        return {}
    except Exception as e:
        logger.error(f"KG: Error processing chunk {chunk_num}: {e}", exc_info=True)
        return {}

def _kg_merge_graphs(graphs: list[dict]) -> dict:
    logger.info("KG: Merging graph fragments...")
    final_nodes = {} # Use dict for easy de-duplication and update
    final_edges = set() # Use set of tuples for de-duplication

    for i, graph_fragment in enumerate(graphs):
        if graph_fragment is None:
            logger.warning(f"KG: Skipping None graph fragment at index {i}.")
            continue
        if not isinstance(graph_fragment, dict) or 'nodes' not in graph_fragment or 'edges' not in graph_fragment:
            logger.warning(f"KG: Skipping invalid graph fragment at index {i}")
            continue

        # Process nodes
        for node in graph_fragment.get('nodes', []):
            if not isinstance(node, dict):
                logger.warning(f"KG: Skipping non-dict node in fragment {i}: {node}")
                continue
            node_id = node.get('id')
            if not node_id or not isinstance(node_id, str): # Ensure ID is a non-empty string
                logger.warning(f"KG: Skipping node with invalid or missing ID in fragment {i}: {node}")
                continue
            
            if node_id not in final_nodes:
                final_nodes[node_id] = node
            else:
                # Merge properties, e.g., longer description, fill missing parent/type
                existing_node = final_nodes[node_id]
                for key in ['description', 'type', 'parent']:
                    if node.get(key) and (not existing_node.get(key) or len(str(node.get(key))) > len(str(existing_node.get(key)))):
                         if key == 'parent' and node.get(key) == node_id: continue # Avoid self-parenting from merge
                         existing_node[key] = node.get(key)


        # Process edges
        for edge in graph_fragment.get('edges', []):
            if not isinstance(edge, dict) or not all(k in edge for k in ['from', 'to', 'relationship']):
                logger.warning(f"KG: Skipping invalid edge (missing keys) in fragment {i}: {edge}")
                continue
            # Ensure all parts of an edge are strings and not empty
            if not all(isinstance(edge.get(k), str) and edge.get(k) for k in ['from', 'to', 'relationship']):
                logger.warning(f"KG: Skipping edge with non-string or empty components in fragment {i}: {edge}")
                continue
            # Avoid self-loops if not desired, or ensure nodes exist
            if edge['from'] == edge['to']: # Optional: disallow self-loops
                # logger.debug(f"KG: Skipping self-loop edge: {edge}")
                continue
            edge_tuple = (edge['from'], edge['to'], edge['relationship'])
            final_edges.add(edge_tuple)

    # Convert back to list format
    merged_graph_nodes = list(final_nodes.values())
    merged_graph_edges = [{"from": e[0], "to": e[1], "relationship": e[2]} for e in final_edges]
    
    # Filter out edges where 'from' or 'to' node doesn't exist in the final_nodes list
    # This can happen if chunks produce partial relationships
    valid_node_ids = set(n['id'] for n in merged_graph_nodes)
    final_merged_edges = [edge for edge in merged_graph_edges if edge['from'] in valid_node_ids and edge['to'] in valid_node_ids]

    merged_graph = {"nodes": merged_graph_nodes, "edges": final_merged_edges}
    logger.info(f"KG: Merged into {len(merged_graph['nodes'])} nodes, {len(merged_graph['edges'])} edges.")
    return merged_graph


def _kg_save_graph(graph: dict, original_doc_filename: str):
    """Saves the graph to a file named after the original document."""
    kg_filepath = get_kg_filepath(original_doc_filename)
    try:
        logger.info(f"KG: Saving graph for '{original_doc_filename}' to {kg_filepath}...")
        os.makedirs(os.path.dirname(kg_filepath), exist_ok=True)
        with open(kg_filepath, 'w', encoding='utf-8') as f:
            json.dump(graph, f, indent=2, ensure_ascii=False)
        logger.info(f"KG: Graph for '{original_doc_filename}' saved to {kg_filepath}.")
    except Exception as e:
        logger.error(f"Failed to save graph to {kg_filepath}: {e}", exc_info=True)


def generate_knowledge_graph_from_pdf(doc_filename: str) -> dict | None: # Renamed from generate_knowledge_graph_from_pdf for clarity
    """
    Generates a knowledge graph from a PDF or PPT file, saving it named after the doc.
    'doc_filename' is the base name of the file, e.g., "MyDocument.pdf".
    """
    global document_texts_cache
    kg_ollama_client = _initialize_kg_ollama_client()
    if not kg_ollama_client:
        logger.error(f"KG: Cannot generate graph for '{doc_filename}'. KG Ollama client not initialized.")
        return None

    logger.info(f"KG: Generating Knowledge Graph for '{doc_filename}'...")
    
    doc_text = document_texts_cache.get(doc_filename)
    if not doc_text:
        logger.debug(f"KG: '{doc_filename}' not in cache. Attempting to load from disk...")
        # Determine full path (check UPLOAD_FOLDER first, then DEFAULT_PDFS_FOLDER)
        file_path_to_load = os.path.join(UPLOAD_FOLDER, doc_filename)
        if not os.path.exists(file_path_to_load):
            file_path_to_load = os.path.join(DEFAULT_PDFS_FOLDER, doc_filename)

        if os.path.exists(file_path_to_load):
            logger.info(f"KG: Found '{doc_filename}' at {file_path_to_load}. Extracting text...")
            doc_text = extract_text_from_file(file_path_to_load)
            if doc_text:
                document_texts_cache[doc_filename] = doc_text # Cache it
                logger.info(f"KG: Cached text for '{doc_filename}' from {file_path_to_load}.")
            else:
                logger.error(f"KG: Failed to extract text from '{doc_filename}' at {file_path_to_load}.")
                return None
        else:
            logger.error(f"KG: File '{doc_filename}' not found in {UPLOAD_FOLDER} or {DEFAULT_PDFS_FOLDER}.")
            return None

    if not doc_text.strip(): # Check if text is not just whitespace
        logger.error(f"KG: No actual text content available for '{doc_filename}' to generate graph.")
        return None

    chunks = _kg_split_into_chunks(doc_text, chunk_size=KG_CHUNK_SIZE, overlap=KG_CHUNK_OVERLAP)
    if not chunks:
        logger.warning(f"KG: No chunks generated for '{doc_filename}'. KG generation aborted.")
        return None

    logger.info(f"KG: Processing {len(chunks)} chunks for '{doc_filename}' using {KG_MAX_WORKERS} workers with model {KG_MODEL}.")
    all_partial_graphs = []
    # Use KG_PROMPT_TEMPLATE from config
    prompt_for_chunks = KG_PROMPT_TEMPLATE 

    with concurrent.futures.ThreadPoolExecutor(max_workers=KG_MAX_WORKERS) as executor:
        # Prepare tasks: tuple of (index, chunk_text)
        tasks = [(i, chunk) for i, chunk in enumerate(chunks)]
        # Submit tasks and map futures back to their original index for logging/debugging
        future_to_index = {
            executor.submit(_kg_process_single_chunk, task_data, kg_ollama_client, KG_MODEL, prompt_for_chunks): task_data[0] 
            for task_data in tasks
        }
        
        for future in tqdm(concurrent.futures.as_completed(future_to_index), total=len(tasks), desc=f"KG: Processing Chunks ({doc_filename})", unit="chunk"):
            original_chunk_index = future_to_index[future]
            try:
                graph_data_fragment = future.result()
                if graph_data_fragment:
                    all_partial_graphs.append(graph_data_fragment)
            except Exception as e: # This catches exceptions from _kg_process_single_chunk if not caught inside
                logger.error(f"KG: Uncaught error from chunk {original_chunk_index + 1} processing for '{doc_filename}': {e}", exc_info=True)

    logger.info(f"KG: Processed {len(all_partial_graphs)} valid graph fragments for '{doc_filename}' from {len(chunks)} chunks.")

    if not all_partial_graphs:
        logger.error(f"KG: No valid graph fragments generated for '{doc_filename}'. Cannot merge.")
        return None

    final_graph = _kg_merge_graphs(all_partial_graphs)
    
    # Save graph using the original document filename to create the KG filename
    _kg_save_graph(final_graph, doc_filename) 

    logger.info(f"KG: Knowledge graph generation for '{doc_filename}' complete.")
    return final_graph


# Main Execution / Test
if __name__ == "__main__":
    # Basic logging for testing
    logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - [%(name)s:%(lineno)d] - %(message)s")
    
    # Setup directories (idempotent)
    for dir_path in [DEFAULT_PDFS_FOLDER, UPLOAD_FOLDER, FAISS_FOLDER, KG_OUTPUT_FOLDER]:
        os.makedirs(dir_path, exist_ok=True)
    
    logger.info("AI Core Testing (with KG Integration)...")
    
    # 1. Initialize AI components (LLM, Embeddings)
    embeddings_instance, llm_instance = initialize_ai_components()
    if not embeddings_instance or not llm_instance:
        logger.error("Failed to initialize core AI components. Aborting test.")
        exit()
    
    # 2. Load vector store (it might be empty initially)
    load_vector_store()

    # 3. Create a dummy PDF for testing (if it doesn't exist)
    dummy_pdf_name = "DevOps-Test-KG.pdf"
    dummy_pdf_path = os.path.join(DEFAULT_PDFS_FOLDER, dummy_pdf_name)
    
    # if not os.path.exists(dummy_pdf_path):
    #     try:
    #         from reportlab.pdfgen import canvas
    #         from reportlab.lib.pagesizes import letter
    #         logger.info(f"Creating dummy PDF: {dummy_pdf_path}")
    #         c = canvas.Canvas(dummy_pdf_path, pagesize=letter)
    #         c.drawString(72, 800, "DevOps Test Document for Knowledge Graph")
    #         c.drawString(72, 780, "This document covers Continuous Integration (CI) and Continuous Deployment (CD).")
    #         c.drawString(72, 760, "Key tools include Jenkins for CI, and Docker for containerization.")
    #         c.drawString(72, 740, "Kubernetes is used for orchestrating Docker containers.")
    #         c.drawString(72, 720, "Monitoring is crucial in DevOps, often using Prometheus.")
    #         c.save()
    #         logger.info(f"Dummy PDF '{dummy_pdf_name}' created.")
    #     except ImportError:
    #         logger.warning("ReportLab not installed. Cannot create dummy PDF. Please create it manually for testing.")
    #     except Exception as e:
    #         logger.error(f"Error creating dummy PDF: {e}")

    if os.path.exists(dummy_pdf_path):
        # 4. Process the dummy PDF: extract text, chunk, add to vector store
        logger.info(f"Processing '{dummy_pdf_name}' for RAG...")
        text_content = extract_text_from_file(dummy_pdf_path)
        if text_content:
            document_texts_cache[dummy_pdf_name] = text_content # Cache it
            chunks = create_chunks_from_text(text_content, dummy_pdf_name)
            if chunks:
                add_documents_to_vector_store(chunks)
                logger.info(f"'{dummy_pdf_name}' processed and added to vector store.")
            else:
                logger.error(f"Could not create chunks for '{dummy_pdf_name}'.")
        else:
            logger.error(f"Could not extract text from '{dummy_pdf_name}'.")

        # 5. Generate Knowledge Graph for the dummy PDF
        logger.info(f"Testing KG generation for {dummy_pdf_name}...")
        kg_result = generate_knowledge_graph_from_pdf(dummy_pdf_name) # Pass base filename
        if kg_result:
            logger.info(f"KG for '{dummy_pdf_name}' generated: {len(kg_result.get('nodes',[]))} nodes, {len(kg_result.get('edges',[]))} edges.")
            kg_file_expected = get_kg_filepath(dummy_pdf_name)
            logger.info(f"KG saved to: {kg_file_expected}")
        else:
            logger.error(f"Failed to generate KG for '{dummy_pdf_name}'.")
    else:
        logger.warning(f"Dummy PDF '{dummy_pdf_name}' not found. Skipping RAG and KG generation tests for it.")

    # 6. Test chat response with RAG and KG integration
    test_query = "What is CI/CD and what tools are used?"
    logger.info(f"\nTesting chat with query: '{test_query}'")
    
    if vector_store: # Ensure RAG can even attempt
        rag_docs, rag_context_text, rag_docs_map = perform_rag_search(test_query)
        logger.info(f"RAG found {len(rag_docs)} documents for the query.")
        logger.debug(f"RAG context text: {rag_context_text[:300]}...")
        
        answer, thinking = synthesize_chat_response(test_query, rag_context_text, rag_docs_map)
        
        logger.info(f"\n--- Test Chat Response ---")
        logger.info(f"Thinking:\n{thinking}")
        logger.info(f"Answer:\n{answer}")
        logger.info(f"--- End Test Chat Response ---")
    else:
        logger.warning("Vector store not available. Skipping chat test.")
        
    logger.info("\nAI Core test (with KG integration) completed.")

# --- END OF FILE ai_core.py ---